"""
Document Extractor Module
Handles text extraction from various file formats
"""

import os
import io
import logging
from typing import Dict, Optional
from pathlib import Path

import PyPDF2
import pdfplumber
from docx import Document
from pptx import Presentation
import openpyxl

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class DocumentExtractor:
    
    SUPPORTED_EXTENSIONS = {
        'pdf': ['.pdf'],
        'word': ['.docx'],
        'powerpoint': ['.pptx'],
        'excel': ['.xlsx'],
        'text': ['.txt', '.md']
    }
    
    def __init__(self):
        """Initialize DocumentExtractor without OCR"""
        logger.info("DocumentExtractor initialized (OCR disabled)")
    
    def is_supported_file(self, filename: str) -> bool:
        ext = Path(filename).suffix.lower()
        return any(ext in extensions for extensions in self.SUPPORTED_EXTENSIONS.values())
    
    def get_file_type(self, filename: str) -> Optional[str]:
        ext = Path(filename).suffix.lower()
        for file_type, extensions in self.SUPPORTED_EXTENSIONS.items():
            if ext in extensions:
                return file_type
        return None
    
    def extract_text(self, file_content: bytes, filename: str) -> Dict:
        file_type = self.get_file_type(filename)
        
        if not file_type:
            return {
                'success': False,
                'text': '',
                'metadata': {},
                'error': f'Unsupported file type: {Path(filename).suffix}'
            }
        
        try:
            if file_type == 'pdf':
                result = self._extract_from_pdf(file_content, filename)
            elif file_type == 'word':
                result = self._extract_from_word(file_content)
            elif file_type == 'powerpoint':
                result = self._extract_from_powerpoint(file_content)
            elif file_type == 'excel':
                result = self._extract_from_excel(file_content)
            elif file_type == 'text':
                result = self._extract_from_text(file_content)
            else:
                result = {
                    'success': False,
                    'text': '',
                    'error': f'Handler not implemented for {file_type}'
                }
            
            if result['success']:
                result['metadata'].update({
                    'filename': filename,
                    'file_type': file_type,
                    'text_length': len(result['text'])
                })
            
            return result
            
        except Exception as e:
            logger.error(f"Error processing {filename}: {str(e)}")
            return {
                'success': False,
                'text': '',
                'metadata': {'filename': filename, 'file_type': file_type},
                'error': str(e)
            }
    
    def _extract_from_pdf(self, file_content: bytes, filename: str) -> Dict:
        text_parts = []
        metadata = {'pages': 0, 'method': 'text_extraction'}
        
        try:
            # Try pdfplumber first (better for complex layouts)
            with pdfplumber.open(io.BytesIO(file_content)) as pdf:
                metadata['pages'] = len(pdf.pages)
                
                for page_num, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text()
                    
                    if page_text and len(page_text.strip()) > 10:
                        text_parts.append(page_text)
                    else:
                        logger.warning(f"Page {page_num} of {filename} has minimal text")
            
            if not text_parts:
                return {
                    'success': False,
                    'text': '',
                    'metadata': metadata,
                    'error': 'No text found in PDF (OCR not available)'
                }
            
            return {
                'success': True,
                'text': '\n\n'.join(text_parts),
                'metadata': metadata,
                'error': None
            }
            
        except Exception as e:
            logger.warning(f"pdfplumber failed for {filename}: {str(e)}, trying PyPDF2")
            
            try:
                # Fallback to PyPDF2
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                metadata['pages'] = len(pdf_reader.pages)
                metadata['method'] = 'PyPDF2'
                
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    page_text = page.extract_text()
                    if page_text and len(page_text.strip()) > 10:
                        text_parts.append(page_text)
                
                if not text_parts:
                    return {
                        'success': False,
                        'text': '',
                        'metadata': metadata,
                        'error': 'No text found in PDF (OCR not available)'
                    }
                
                return {
                    'success': True,
                    'text': '\n\n'.join(text_parts),
                    'metadata': metadata,
                    'error': None
                }
            except Exception as pdf_error:
                return {
                    'success': False,
                    'text': '',
                    'metadata': metadata,
                    'error': f'PDF processing failed: {str(pdf_error)}'
                }
    
    def _extract_from_word(self, file_content: bytes) -> Dict:
        try:
            doc = Document(io.BytesIO(file_content))
            
            text_parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(row_text)
            
            metadata = {
                'paragraphs': len(doc.paragraphs),
                'tables': len(doc.tables)
            }
            
            return {
                'success': True,
                'text': '\n\n'.join(text_parts),
                'metadata': metadata,
                'error': None
            }
            
        except Exception as e:
            return {
                'success': False,
                'text': '',
                'metadata': {},
                'error': f'Word document processing failed: {str(e)}'
            }
    
    def _extract_from_powerpoint(self, file_content: bytes) -> Dict:
        try:
            prs = Presentation(io.BytesIO(file_content))
            
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)
                    
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                            if row_text.strip():
                                slide_texts.append(row_text)
                
                if slide_texts:
                    text_parts.append(f"[Slide {slide_num}]\n" + '\n'.join(slide_texts))
            
            metadata = {'slides': len(prs.slides)}
            
            return {
                'success': True,
                'text': '\n\n'.join(text_parts),
                'metadata': metadata,
                'error': None
            }
            
        except Exception as e:
            return {
                'success': False,
                'text': '',
                'metadata': {},
                'error': f'PowerPoint processing failed: {str(e)}'
            }
    
    def _extract_from_excel(self, file_content: bytes) -> Dict:
        try:
            workbook = openpyxl.load_workbook(io.BytesIO(file_content), data_only=True)
            
            text_parts = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"[Sheet: {sheet_name}]")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                    if row_text.strip():
                        text_parts.append(row_text)
            
            metadata = {'sheets': len(workbook.sheetnames)}
            
            return {
                'success': True,
                'text': '\n'.join(text_parts),
                'metadata': metadata,
                'error': None
            }
            
        except Exception as e:
            return {
                'success': False,
                'text': '',
                'metadata': {},
                'error': f'Excel processing failed: {str(e)}'
            }
    
    def _extract_from_text(self, file_content: bytes) -> Dict:
        try:
            encodings = ['utf-8', 'latin-1', 'cp1252']
            text = None
            used_encoding = None
            
            for encoding in encodings:
                try:
                    text = file_content.decode(encoding)
                    used_encoding = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            if text is None:
                raise ValueError("Could not decode text file")
            
            return {
                'success': True,
                'text': text,
                'metadata': {'encoding': used_encoding},
                'error': None
            }
            
        except Exception as e:
            return {
                'success': False,
                'text': '',
                'metadata': {},
                'error': f'Text file processing failed: {str(e)}'
            }
